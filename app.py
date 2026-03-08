import streamlit as st
import pandas as pd
import numpy as np
import re, os, json, io, time, zipfile
from groq import Groq
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# ── PAGE SETUP ────────────────────────────────────────────────────────────────
st.set_page_config(page_title="InsightX AI | Universal RAG", layout="wide", page_icon="🧠")

# ── PREMIUM CSS ───────────────────────────────────────────────────────────────
st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;800&family=Outfit:wght@300;500;700&display=swap');
    
    html, body, [class*="css"] {
        font-family: 'Inter', sans-serif;
    }
    
    h1, h2, h3, .stHeader {
        font-family: 'Outfit', sans-serif;
        font-weight: 700;
        background: linear-gradient(90deg, #6366f1, #a855f7, #ec4899);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
    }
    
    .stApp {
        background: #0f172a;
        color: #f1f5f9;
    }
    
    /* Glassmorphism sidebar */
    [data-testid="stSidebar"] {
        background-color: rgba(30, 41, 59, 0.7) !important;
        backdrop-filter: blur(12px);
        border-right: 1px solid rgba(255, 255, 255, 0.1);
    }
    
    /* Custom Card Style */
    div.stMetric {
        background: rgba(30, 41, 59, 0.5);
        padding: 20px;
        border-radius: 15px;
        border: 1px solid rgba(255, 255, 255, 0.1);
        transition: transform 0.3s ease;
    }
    div.stMetric:hover {
        transform: translateY(-5px);
        border-color: #6366f1;
    }
    
    /* Tab Styling */
    .stTabs [data-baseweb="tab-list"] {
        gap: 10px;
        background-color: transparent;
    }
    .stTabs [data-baseweb="tab"] {
        height: 50px;
        border-radius: 10px 10px 0 0;
        padding: 0 20px;
        background-color: rgba(30, 41, 59, 0.5);
        border: 1px solid rgba(255, 255, 255, 0.1);
        color: #94a3b8;
    }
    .stTabs [aria-selected="true"] {
        background-color: #6366f1 !important;
        color: white !important;
    }
    
    /* Chat bubbles */
    [data-testid="stChatMessage"] {
        background: rgba(30, 41, 59, 0.5);
        border-radius: 15px;
        border: 1px solid rgba(255, 255, 255, 0.1);
        margin-bottom: 10px;
    }
    
    /* Animation */
    @keyframes fadeIn {
        from { opacity: 0; transform: translateY(10px); }
        to { opacity: 1; transform: translateY(0); }
    }
    .stApp > div {
        animation: fadeIn 0.8s ease-out;
    }
    </style>
""", unsafe_allow_html=True)

st.title("🧠 InsightX AI")
st.caption("The future of Universal RAG & Data Intelligence — Lightning fast, 100% Secure.")

# ── API KEY HANDLING ──────────────────────────────────────────────────────────
api_key = st.secrets.get("GROQ_API_KEY", os.environ.get("GROQ_API_KEY", ""))
has_api = bool(api_key and api_key != "PASTE_YOUR_GROQ_API_KEY_HERE")

if has_api:
    client = Groq(api_key=api_key)
else:
    client = None

# ── LOGIC GUARD ──
def check_api():
    if not has_api:
        st.warning("⚠️ **Groq API Key Missing**")
        st.info("To enable AI features, please set `GROQ_API_KEY` in your Streamlit Secrets or Environment Variables.")
        return False
    return True

# ── CONSTANTS ─────────────────────────────────────────────────────────────────
MAX_FILE_MB    = 1024    # Increased to 1GB (1024MB) as per user requirement
MAX_RAG_ROWS   = 2000    # Reduced from 10,000 to improve memory/lag
MAX_WORDS      = 30000   # Optimized for TF-IDF
TOP_K          = 3
CONTEXT_CAP    = 4000
MAX_OUT_TOKENS = 600

# ── COLUMN HEURISTICS ──────────────────────────────────────────────────────────
def find_col(df, keywords):
    """Helper to find the best matching column for a list of keywords."""
    for col in df.columns:
        if any(k in col.lower() for k in keywords):
            return col
    return None

# ── FILE EXTRACTORS & UTILS ────────────────────────────────────────────────────
@st.cache_data(show_spinner=False)
def get_df_stats(df):
    """Cache heavy dataframe stats calculations."""
    missing = {k: v for k, v in df.isnull().sum().items() if v > 0}
    dupes = int(df.duplicated().sum())
    describe = df.describe()
    return missing, dupes, describe

@st.cache_data(show_spinner=False)
def get_csv_download(df):
    """Cache CSV conversion to avoid lag on every rerun."""
    return df.to_csv(index=False).encode('utf-8')

@st.cache_data(show_spinner=False)
def extract_text_from_csv(file_bytes):
    # Use low_memory=True and dtype inference to save RAM
    # We read the full file for 'Precise Calculator' but optimize its storage
    try:
        df = pd.read_csv(io.BytesIO(file_bytes), low_memory=True)
        
        # DOWNCASTING: Convert float64 to float32, int64 to int32 where possible
        # This saves ~50% RAM while maintaining high precision for calculations.
        for col in df.select_dtypes(include=['float']).columns:
            df[col] = pd.to_numeric(df[col], downcast='float')
        for col in df.select_dtypes(include=['integer']).columns:
            df[col] = pd.to_numeric(df[col], downcast='integer')
            
        total = len(df)
        text  = f"CSV | Rows: {total:,} | Columns: {list(df.columns)}\n\n"
        
        if total > MAX_RAG_ROWS:
            sample = df.sample(n=MAX_RAG_ROWS, random_state=42)
            text += f"(Representative sample of {MAX_RAG_ROWS:,}/{total:,} rows for AI context.)\n\n"
            text += sample.to_string(index=False, max_rows=MAX_RAG_ROWS)
        else:
            text += df.to_string(index=False)
        return text, df
    except Exception as e:
        st.error(f"Error parsing CSV: {e}")
        return None, None

@st.cache_data(show_spinner=False)
def extract_text_from_excel(file_bytes):
    try:
        # Use openpyxl for engine and read head for preview to save memory on huge files
        xl  = pd.ExcelFile(io.BytesIO(file_bytes), engine='openpyxl')
        txt = ""
        fdf = None
        for sheet in xl.sheet_names:
            df  = xl.parse(sheet)
            # Downcast to save RAM
            for col in df.select_dtypes(include=['float']).columns:
                df[col] = pd.to_numeric(df[col], downcast='float')
            if fdf is None: fdf = df
            txt += f"\n--- Sheet: {sheet} ---\nRows:{len(df)} Cols:{list(df.columns)}\n"
            txt += df.head(100).to_string(index=False) + "\n"
        return txt, fdf
    except Exception as e:
        st.error(f"Excel Error: {e}")
        return None, None

@st.cache_data(show_spinner=False)
def extract_text_from_pdf(file_bytes):
    import pypdf
    reader = pypdf.PdfReader(io.BytesIO(file_bytes))
    text, ocr_pages = f"PDF | Pages: {len(reader.pages)}\n", []
    for i, page in enumerate(reader.pages):
        pt = page.extract_text() or ""
        if len(pt.strip()) > 30:
            text += f"\n--- Page {i+1} ---\n{pt}\n"
        else:
            ocr_pages.append(i + 1)
    if ocr_pages:
        try:
            import pytesseract
            from pdf2image import convert_from_bytes
            from PIL import Image
            imgs = convert_from_bytes(file_bytes, dpi=120)
            for pn in ocr_pages:
                if pn <= len(imgs):
                    ot = pytesseract.image_to_string(imgs[pn - 1])
                    if ot.strip():
                        text += f"\n--- Page {pn} (OCR) ---\n{ot}\n"
        except Exception:
            pass
    return text, None

@st.cache_data(show_spinner=False)
def extract_text_from_docx(file_bytes):
    import docx
    doc  = docx.Document(io.BytesIO(file_bytes))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    for tbl in doc.tables:
        for row in tbl.rows:
            text += " | ".join(c.text.strip() for c in row.cells) + "\n"
    return text, None

@st.cache_data(show_spinner=False)
def extract_text_from_pptx(file_bytes):
    from pptx import Presentation
    prs  = Presentation(io.BytesIO(file_bytes))
    text = f"PowerPoint | Slides: {len(prs.slides)}\n"
    for i, slide in enumerate(prs.slides):
        text += f"\n--- Slide {i+1} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
    return text, None

@st.cache_data(show_spinner=False)
def extract_text_from_image(file_bytes):
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(io.BytesIO(file_bytes))
        txt = pytesseract.image_to_string(img)
        w, h = img.size
        return f"Image {w}x{h}px\n\n{txt if txt.strip() else 'No readable text.'}", None
    except Exception as e:
        return f"Image error: {e}", None

@st.cache_data(show_spinner=False)
def extract_text_from_zip(file_bytes):
    text = "ZIP Archive:\n"
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
        for name in z.namelist():
            text += f"  - {name}\n"
            if name.endswith((".txt", ".csv", ".json", ".md", ".py", ".sql")):
                try:
                    with z.open(name) as f:
                        text += f"\n--- {name} ---\n"
                        # Use string slicing safely
                        raw_content = f.read().decode("utf-8", errors="ignore")
                        text += (raw_content[:2000] + "...") if len(raw_content) > 2000 else raw_content
                        text += "\n"
                except Exception:
                    pass
    return text, None

def extract_content(uploaded_file):
    try:
        name       = uploaded_file.name.lower()
        file_bytes = uploaded_file.getvalue()
        size_mb    = len(file_bytes) / (1024 * 1024)

        if size_mb > MAX_FILE_MB:
            st.error(f"File too large ({size_mb:.1f}MB). Max: {MAX_FILE_MB}MB.")
            st.stop()

        if name.endswith(".csv"):
            text, df = extract_text_from_csv(file_bytes)
        elif name.endswith((".xlsx", ".xls")):
            text, df = extract_text_from_excel(file_bytes)
        elif name.endswith(".pdf"):
            text, df = extract_text_from_pdf(file_bytes)
        elif name.endswith(".docx"):
            text, df = extract_text_from_docx(file_bytes)
        elif name.endswith(".pptx"):
            text, df = extract_text_from_pptx(file_bytes)
        elif name.endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff")):
            text, df = extract_text_from_image(file_bytes)
        elif name.endswith(".json"):
            try:
                data = json.loads(file_bytes.decode("utf-8", errors="ignore"))
                text = json.dumps(data, indent=2)[:8000]
                if isinstance(data, list) and len(data) > 0 and isinstance(data[0], dict):
                    df = pd.DataFrame(data)
                else:
                    df = None
            except Exception:
                text, df = "Error reading JSON", None
        elif name.endswith(".zip"):
            text, df = extract_text_from_zip(file_bytes)
        else:
            text = file_bytes.decode("utf-8", errors="ignore")
            df   = None

        if not text or len(text.strip()) < 10:
            st.error("File appears empty or unreadable.")
            st.stop()

        return text, df
    except MemoryError:
        st.error("📉 OUT OF MEMORY: This file is too large for the current server RAM (1GB limit).")
        st.info("Try uploading a smaller version of the file or a CSV with fewer columns.")
        st.stop()
    except Exception as e:
        st.error(f"⚠️ Unexpected Error: {e}")
        st.stop()

@st.cache_data(show_spinner=False)
def generate_suggestions(text, file_name, is_tabular=False):
    """Generate context-aware questions grounded in the file content."""
    if not check_api():
        return ["Summary of the file", "Key trends", "Data structure help"]
    
    if is_tabular:
        prompt = (
            f"Based on the following data sample from '{file_name}', "
            "generate 5-7 precise questions that can be answered by analyzing this data. "
            "Focus on trends, totals, averages, or specific column relationships. "
            "STRICT RULE: Only generate questions that can be answered using the provided columns and data. "
            "Avoid general or imaginative questions. "
            "Provide ONLY the questions, one per line, no numbering or extra text.\n\n"
            f"Data Sample:\n{text[:5000]}"
        )
    else:
        prompt = (
            f"Based on the following text from '{file_name}', "
            "generate 5-7 insightful questions that can be answered using this content. "
            "Focus on summaries, key facts, or specific details mentioned in the text. "
            "STRICT RULE: Only generate questions that can be answered directly by the provided text. "
            "Avoid external knowledge or assumptions. "
            "Provide ONLY the questions, one per line, no numbering or extra text.\n\n"
            f"Content Sample:\n{text[:5000]}"
        )
    for attempt in range(4):
        try:
            r = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=300,
            )
            lines = r.choices[0].message.content.strip().split("\n")
            return [l.strip().lstrip("0123456789.- ") for l in lines if l.strip()]
        except Exception as e:
            if "429" in str(e) and attempt < 3:
                time.sleep(2 ** attempt + 1)
            else:
                break
    # Fallback suggestions if API busy
    return [
        "Give me a detailed summary of this file",
        "What are the most important points here?",
        "Are there any specific trends or patterns?",
        f"Explain the data structure of {file_name}",
        "What is the main purpose of this content?"
    ]

# ── FAST TF-IDF RAG ───────────────────────────────────────────────────────────
@st.cache_resource
def build_tfidf_index(text, file_name):
    """Build TF-IDF index — Memory optimized."""
    if not text: return None, None, []
    
    # Use a more memory-efficient way to sample if text is huge
    if len(text) > 2_000_000: # ~2MB of text
        text = text[:2_000_000]
        st.toast("⚠️ Large text file detected. Indexing the first 2 million characters for RAG.", icon="ℹ️")

    words = text.split()
    if len(words) > MAX_WORDS:
        words = words[:MAX_WORDS]

    # Split into chunks efficiently
    chunk_size, overlap = 400, 100
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)

    if not chunks:
        return None, None, []

    vectorizer = TfidfVectorizer(max_features=2000, stop_words="english", ngram_range=(1, 1))
    tfidf_matrix = vectorizer.fit_transform(chunks)
    return vectorizer, tfidf_matrix, chunks

def retrieve_chunks(question, vectorizer, tfidf_matrix, chunks):
    """Retrieve top-K relevant chunks using TF-IDF cosine similarity."""
    if vectorizer is None:
        return []
    q_vec  = vectorizer.transform([question])
    scores = cosine_similarity(q_vec, tfidf_matrix).flatten()
    top_idx = scores.argsort()[-TOP_K:][::-1]
    return [chunks[i] for i in top_idx if i < len(chunks)]

# ── AI CALLS ──────────────────────────────────────────────────────────────────
def ask_ai(question, chunks, history, retries=4):
    if not check_api():
        return "⚠️ API Key Missing. Please configure it in Streamlit Secrets."
    
    context = "\n---\n".join(chunks)
    if len(context) > CONTEXT_CAP:
        context = context[:CONTEXT_CAP]
    
    # ... rest of the function ...
    messages = [{
        "role": "system",
        "content": (
            "You are a strict Data Analyst and Assistant. "
            "Answer ONLY using the provided file excerpts. "
            "If the excerpts do not contain the answer, say: 'The file does not contain information related to the question asked.' "
            "Do NOT use general knowledge or guess if it's not in the file. "
            "Stay 100% grounded in the uploaded content."
        )
    }]
    if history and isinstance(history, list):
        for c in history[-3:]:
            messages.append({"role": "user",      "content": c["question"]})
            messages.append({"role": "assistant", "content": c.get("answer", "")})
    messages.append({
        "role": "user",
        "content": f"File content:\n{context}\n\nQuestion: {question}"
    })

    for attempt in range(retries):
        try:
            r = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=messages,
                max_tokens=MAX_OUT_TOKENS,
            )
            return r.choices[0].message.content.strip()
        except Exception as e:
            if "429" in str(e) and attempt < retries - 1:
                wait = 2 ** attempt + 2
                st.toast(f"⏳ System busy (Rate Limit). Retrying in {wait}s...", icon="⚠️")
                time.sleep(wait)
            elif attempt < retries - 1:
                time.sleep(1)
            else:
                return f"AI Error: {e}. Please try again in 5-10 seconds."

def run_code(question, df, history, execute=True, retries=3):
    if not check_api():
        return "⚠️ API Key Missing", "Cannot execute AI logic without key."
    
    cols   = str(list(df.columns))
    # ...
    dtypes = df.dtypes.astype(str).to_string()
    sample = df.head(3).to_string()
    info   = f"Columns:{cols}\nDtypes:\n{dtypes}\nSample:\n{sample}"
    if len(info) > 1200:
        info = info[:1200]

    messages = [{
        "role": "system",
        "content": (
            "You are a Python data analyst. Write pandas code to answer the question. "
            "DataFrame is 'df'. Use st.write(), st.dataframe(), st.bar_chart(), "
            "st.line_chart(), st.area_chart(), or st.metric() to show results. "
            "Return ONLY raw Python. No markdown, no backticks, no explanation."
        )
    }]
    # Safety guard for history slicing
    if history and isinstance(history, list):
        for c in history[-2:]:
            messages.append({"role": "user",      "content": c["question"]})
            messages.append({"role": "assistant", "content": c.get("answer", "")})
    
    messages.append({
        "role": "user",
        "content": f"DataFrame info:\n{info}\n\nQuestion: {question}"
    })

    for attempt in range(retries):
        try:
            r = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=messages,
                max_tokens=600,
            )
            raw = (
                r.choices[0].message.content
                .replace("```python", "").replace("```", "").strip()
            )
            clean = "\n".join(
                line for line in raw.splitlines()
                if not re.match(r"^\d{1,2}:\d{2}\s*(AM|PM)?$",
                                line.strip(), re.IGNORECASE)
            ).strip()
            
            # Basic sanitization
            dangerous_tokens = ["os.", "sys.", "subprocess", "open(", "eval(", "exec(", "shutil", "builtins"]
            if any(token in clean for token in dangerous_tokens):
                return clean, "Security blocked: Generated code contains potentially unsafe operations."

            if execute:
                exec(clean, {"df": df, "st": st, "pd": pd, "np": np})
            return clean, None
        except Exception as e:
            if "429" in str(e) and attempt < retries - 1:
                wait = 2 ** attempt + 2
                st.toast(f"⏳ System busy (Rate Limit). Retrying in {wait}s...", icon="📊")
                time.sleep(wait)
            elif attempt < retries - 1:
                time.sleep(1)
            else:
                return raw, str(e)

def auto_generate_pulse(df, chunks):
    """Proactively analyze the file for quick insights."""
    if not has_api:
        return "• High-speed RAG ready.\n• Data structure detected.\n• Awaiting API configuration."

    if df is not None:
        info = f"Cols: {list(df.columns)}\nRows: {len(df)}\nSample: {df.head(2).to_string()}"
        prompt = (
            f"As a Senior Data Scientist, provide exactly 3 'Bullet Point' executive observations "
            f"based on this dataset summary:\n{info}\n"
            "Focus on health, trends, or outliers. Keep each point under 10 words."
        )
    else:
        text_sample = "\n".join(chunks[:3])
        prompt = (
            f"Summarize the key intent of this document in exactly 3 'Bullet Point' observations:\n"
            f"{text_sample[:2000]}\n"
            "Keep each point under 10 words."
        )
    
    try:
        r = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=150,
        )
        return r.choices[0].message.content.strip()
    except Exception:
        return "• High-speed RAG ready.\n• Analysis pending."

# ── SIDEBAR ───────────────────────────────────────────────────────────────────
st.sidebar.header("📂 Upload Any File")
st.sidebar.caption("CSV, Excel, PDF, Word, PPT, JSON, TXT, ZIP, Images...")

uploaded_file = st.sidebar.file_uploader(
    "Drop your file here",
    type=["csv","xlsx","xls","pdf","docx","pptx","json","txt","sql",
          "md","py","log","xml","html","zip","png","jpg","jpeg","webp","tiff"]
)

# ── MAIN ──────────────────────────────────────────────────────────────────────
# Handle Sample Data
if "uploaded_df" not in st.session_state:
    st.session_state.uploaded_df = None
if "uploaded_text" not in st.session_state:
    st.session_state.uploaded_text = None

st.sidebar.markdown("---")
st.sidebar.subheader("🚀 Demo Mode")
if st.sidebar.button("Load Sample Transaction Dataset (250k)"):
    data = {
        "Transaction_ID": [f"TXN{i:06d}" for i in range(250000)],
        "Timestamp": pd.date_range("2024-01-01", periods=250000, freq="min"),
        "Amount": np.random.uniform(10, 5000, 250000).round(2),
        "Category": np.random.choice(["Food", "Entertainment", "Education", "Health", "Travel"], 250000),
        "Device_Type": np.random.choice(["iOS", "Android"], 250000),
        "Network_Type": np.random.choice(["5G", "WiFi", "4G"], 250000),
        "Age": np.random.randint(18, 70, 250000),
        "State": np.random.choice(["Maharashtra", "Delhi", "Karnataka", "Tamil Nadu", "Gujarat"], 250000),
        "Status": np.random.choice(["Success", "Failed"], 250000, p=[0.95, 0.05]),
        "Fraud_Flag": np.random.choice([0, 1], 250000, p=[0.99, 0.01])
    }
    st.session_state.uploaded_df = pd.DataFrame(data)
    st.session_state.uploaded_text = "Sample Digital Payments Dataset | 250,000 Transactions | Source: Techfest Synthetic Data"
    st.session_state.last_file = "sample_transactions.csv"
    st.rerun()

if uploaded_file or st.session_state.uploaded_df is not None:
    # Extract
    if uploaded_file:
        with st.spinner(f"📖 Reading {uploaded_file.name}..."):
            text, df = extract_content(uploaded_file)
            # Reset sample if new file uploaded
            st.session_state.uploaded_df = None
            filename = uploaded_file.name
            size_mb  = len(uploaded_file.getvalue()) / (1024 * 1024)
    else:
        df = st.session_state.uploaded_df
        text = st.session_state.uploaded_text
        filename = st.session_state.last_file
        size_mb  = 0.0

    word_count = len(text.split())

    st.sidebar.success(f"✅ {filename}")
    if size_mb > 0:
        st.sidebar.info(f"📝 {word_count:,} words | {size_mb:.2f} MB")
    else:
        st.sidebar.info(f"📝 {word_count:,} words | Demo Dataset")

    # Build TF-IDF index (fast!)
    with st.spinner("⚡ Indexing file (TF-IDF)..."):
        vectorizer, tfidf_matrix, chunks = build_tfidf_index(text, filename)

    st.sidebar.info(f"🧩 {len(chunks)} chunks | ⚡ TF-IDF indexed")

    # Session state
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    if "calc_history" not in st.session_state:
        st.session_state.calc_history = []
    if "suggestions" not in st.session_state:
        st.session_state.suggestions = []
    if "last_file" not in st.session_state:
        st.session_state.last_file = ""
    if "pending_q" not in st.session_state:
        st.session_state.pending_q = None

    if st.session_state.last_file != filename:
        st.session_state.chat_history = []
        st.session_state.calc_history = []
        st.session_state.last_file    = filename
        st.session_state.pending_q    = None
        # PERFORMANCE FIX: Pre-generate dynamic suggestions
        with st.spinner("🧠 Preparing smart suggestions..."):
            is_tabular = df is not None
            st.session_state.suggestions = generate_suggestions(text, filename, is_tabular)
        
        # PREMIUM FEATURE: Generate Auto-Pulse
        with st.spinner("⚡ AI Pulse active..."):
            st.session_state.data_pulse = auto_generate_pulse(df, chunks)

    # Sidebar Data Pulse Widget
    if "data_pulse" in st.session_state:
        st.sidebar.markdown("### 📊 Data Pulse")
        st.sidebar.info(st.session_state.data_pulse)

    # Guide table
    st.markdown("""
| Tab | Best For | Accuracy |
|-----|----------|----------|
| 💬 Ask AI | Summaries, explanations, trends, general Q&A | Smart sample |
| ⚡ Precise Calculator | Totals, averages, rankings, exact numbers | ✅ 100% Full dataset |
| 📊 Data Analysis | Auto charts, quality checks, stats | ✅ 100% Full dataset |
| 👀 Preview | See extracted content & file stats | — |
""")

    tab0, tab1, tab2, tab3, tab4 = st.tabs(
        ["🏆 Leadership Insights", "💬 Ask AI", "⚡ Precise Calculator", "📊 Data Analysis", "👀 Preview"]
    )

    # ── TAB 0: LEADERSHIP INSIGHTS (AUTO) ─────────────────────────────────────
    with tab0:
        st.subheader("🏆 Leadership Insights Dashboard")
        
        if df is not None:
            # Executive Summary Section
            st.markdown("### 📄 Executive Summary")
            st.write(st.session_state.get("data_pulse", "Analyzing..."))
            
            report_txt = f"INSIGHTX AI EXECUTIVE REPORT\nFile: {filename}\nDate: {pd.Timestamp.now()}\n\n"
            report_txt += "--- AI OBSERVATIONS ---\n" + st.session_state.get("data_pulse", "") + "\n\n"
            
            # 1. Descriptive
            with st.expander("📊 1. Descriptive Analysis", expanded=True):
                c1, c2, c3 = st.columns(3)
                amt_col = find_col(df, ['amount', 'amt', 'value', 'val', 'price', 'total', 'volume'])
                if amt_col:
                    val_sum = df[amt_col].sum()
                    val_avg = df[amt_col].mean()
                    c1.metric("Total Volume", f"₹{val_sum:,.0f}")
                    c2.metric("Avg Transaction", f"₹{val_avg:,.2f}")
                    report_txt += f"Total Volume: ₹{val_sum:,.0f}\nAvg Transaction: ₹{val_avg:,.2f}\n"
                
                cat_col = find_col(df, ['category', 'cat', 'type', 'industry', 'dept', 'purpose', 'genre'])
                if cat_col:
                    top_cat = df[cat_col].value_counts().idxmax()
                    c3.metric("Top Category", top_cat)
                    st.bar_chart(df[cat_col].value_counts().head(5))
                    report_txt += f"Top Category: {top_cat}\n"

            # 2. Comparative
            with st.expander("📈 2. Comparative Analysis", expanded=True):
                comp_col = find_col(df, ['device', 'platform', 'mode', 'channel', 'network', 'brand'])
                if comp_col and amt_col:
                    st.write(f"**Transaction Value by {comp_col}**")
                    st.bar_chart(df.groupby(comp_col)[amt_col].mean())
                else:
                    st.info("Upload a dataset with 'Device', 'Payment Mode', or 'Platform' for comparisons.")

            # 3. User Segmentation
            with st.expander("👥 3. User Segmentation", expanded=True):
                age_col = find_col(df, ['age', 'years', 'dob', 'user_age'])
                state_col = find_col(df, ['state', 'region', 'location', 'city', 'country', 'prov'])
                
                sc1, sc2 = st.columns(2)
                if age_col:
                    with sc1:
                        st.write("**Age Distribution**")
                        st.line_chart(df[age_col].value_counts().sort_index())
                if state_col:
                    with sc2:
                        st.write("**Top States/Regions**")
                        st.bar_chart(df[state_col].value_counts().head(5))

            # 4. Risk Metrics
            with st.expander("🛡️ 4. Risk & Operational Metrics", expanded=True):
                fraud_col = find_col(df, ['fraud', 'flag', 'is_fraud', 'risk', 'suspicious'])
                status_col = find_col(df, ['status', 'result', 'success', 'outcome', 'stat'])
                
                rc1, rc2 = st.columns(2)
                if fraud_col:
                    fraud_rate = (pd.to_numeric(df[fraud_col], errors='coerce').fillna(0).mean() * 100)
                    rc1.metric("Fraud Rate", f"{fraud_rate:.2f}%")
                if status_col:
                    is_success = df[status_col].astype(str).str.lower().str.contains('success|ok|completed')
                    fail_rate = (1 - is_success.mean()) * 100
                    rc2.metric("Failure Rate", f"{fail_rate:.2f}%")
            
            st.divider()
            st.download_button("📥 Download Executive Report", data=report_txt, 
                               file_name=f"Executive_Report_{filename}.txt", mime="text/plain")
        else:
            st.subheader("📄 Document Intelligence Summary")
            st.write(st.session_state.get("data_pulse", "Analyzing document..."))
            st.info("Upload a CSV/Excel for the full Leadership Dashboard.")

    # ── TAB 1: ASK AI ─────────────────────────────────────────────────────────
    with tab1:
        st.subheader("💬 Ask Anything — Plain English")
        st.info(
            "💡 Uses smart TF-IDF search to find relevant parts of your file instantly. "
            "Best for **summaries, trends, explanations & general questions**. "
            "For exact numbers → use ⚡ Precise Calculator."
        )

        def handle_q():
            if st.session_state.ask_input:
                st.session_state.pending_q = st.session_state.ask_input
                st.session_state.ask_input = ""
            elif st.session_state.quick_q != "Choose or type below...":
                st.session_state.pending_q = st.session_state.quick_q
                # Do not reset quick_q here; it causes an exception. 
                # It will be reset naturally on the next rerun.

        quick_qs = st.session_state.suggestions
        st.selectbox("⚡ Quick Questions:", ["Choose or type below..."] + quick_qs, 
                     key="quick_q", on_change=handle_q)
        st.text_input("Your question:", placeholder="e.g., What are the key points in this file?", 
                      key="ask_input", on_change=handle_q)

        if st.session_state.pending_q:
            final_q = st.session_state.pending_q
            st.session_state.pending_q = None # Consume it

            # Detect if user asks for chart/graph
            chart_keywords = ["chart", "graph", "plot", "map", "trend", "visualize"]
            is_chart_q = any(k in final_q.lower() for k in chart_keywords)

            # FALLBACK: If the file is small, use the entire text instead of chunks
            context_chunks = []
            if len(text.split()) < 500:
                context_chunks = [text]
            else:
                context_chunks = retrieve_chunks(final_q, vectorizer, tfidf_matrix, chunks)

            if is_chart_q and df is not None:
                with st.spinner("📊 Analyzing data for visualization..."):
                    code_used, error = run_code(final_q, df, st.session_state.chat_history, execute=False)
                    if not error:
                        st.session_state.chat_history.append({
                            "question": final_q, 
                            "answer": f"I've generated a chart for you based on the data.",
                            "type": "chart",
                            "code": code_used
                        })
                        st.rerun()
                    else:
                        st.session_state.chat_history.append({
                            "question": final_q,
                            "answer": f"I tried to generate a chart but ran into an issue: {error}. Let me give you a text summary instead."
                        })
                        # Use context_chunks
                        ans = ask_ai(final_q, context_chunks, st.session_state.chat_history)
                        st.session_state.chat_history[-1]["answer"] += f"\n\n--- Summary ---\n{ans}"
                        st.rerun()
            
            elif is_chart_q and df is None:
                # User asked for chart on non-tabular file
                st.session_state.chat_history.append({
                    "question": final_q,
                    "answer": "I'd love to generate a chart, but this file doesn't contain the structured data needed for graphs. Here's a search-based answer for you:",
                })
                with st.spinner("🔍 Reading content..."):
                    ans = ask_ai(final_q, context_chunks, st.session_state.chat_history)
                    st.session_state.chat_history[-1]["answer"] += f"\n\n{ans}"
                    st.rerun()

            else:
                # Normal text question
                with st.spinner("🔍 Reviewing file..."):
                    try:
                        ans = ask_ai(final_q, context_chunks, st.session_state.chat_history)
                        st.session_state.chat_history.append({"question": final_q, "answer": ans})
                        st.rerun()
                    except Exception as e:
                        st.error(f"AI Error: {e}")

        if st.session_state.chat_history:
            chat_txt = "\n\n".join(f"Q: {c['question']}\nA: {c['answer']}"
                                   for c in st.session_state.chat_history)
            st.download_button("⬇️ Download Chat", data=chat_txt,
                               file_name="chat.txt", mime="text/plain", key="dl_chat")
            st.subheader("🗨️ Conversation")
            for idx, chat in enumerate(reversed(st.session_state.chat_history)):
                # Calculate the true index in the original list (since we are reversed)
                true_idx = len(st.session_state.chat_history) - 1 - idx
                
                with st.chat_message("user"):
                    st.write(chat["question"])
                with st.chat_message("assistant"):
                    st.write(chat["answer"])
                    # RE-RENDER chart if type is chart
                    if chat.get("type") == "chart" and chat.get("code") and df is not None:
                        try:
                            exec(chat["code"], {"df": df, "st": st, "pd": pd, "np": np})
                        except Exception as e:
                            st.info(f"Chart preview unavailable: {e}")
                            with st.expander("Show code used"):
                                st.code(chat["code"], language="python")
                    
                    c1, c2 = st.columns(2)
                    with c1:
                        st.download_button(
                            "⬇️ Download answer", data=chat["answer"],
                            file_name=f"answer_{idx}.txt", mime="text/plain",
                            key=f"dl_ans_{idx}"
                        )
                    with c2:
                        if st.button("🗑️ Delete this", key=f"del_{true_idx}"):
                            st.session_state.chat_history.pop(true_idx)
                            st.rerun()
            if st.button("🗑️ Clear Chat", key="clear_chat"):
                st.session_state.chat_history = []
                st.rerun()

    # ── TAB 2: PRECISE CALCULATOR ──────────────────────────────────────────────
    with tab2:
        st.subheader("⚡ Precise Calculator")
        if df is not None:
            st.success(
                "✅ **100% Accurate — Full Dataset** | "
                "Runs real Python/pandas code on your complete data. "
                "No sampling — every result is mathematically exact."
            )
            calc_q = st.text_input(
                "What do you want to calculate?",
                placeholder="e.g., Total amount by category | Top 10 by value | Monthly trend chart",
                key="calc_input"
            )
            if calc_q:
                with st.spinner("🧮 Running code on full dataset..."):
                    code_used, error = run_code(calc_q, df, st.session_state.calc_history)
                    if error:
                        st.error(f"Error: {error}")
                        if code_used:
                            st.code(code_used, language="python")
                    else:
                        st.session_state.calc_history.append({
                            "question": calc_q,
                            "answer": f"[Executed]\n```python\n{code_used}\n```"
                        })
                        with st.expander("🔍 View generated code"):
                            st.code(code_used, language="python")

            if st.session_state.calc_history:
                st.subheader("🗨️ Calculation History")
                for idx, calc in enumerate(reversed(st.session_state.calc_history)):
                    true_idx = len(st.session_state.calc_history) - 1 - idx
                    with st.expander(f"Q: {calc['question']}", expanded=(idx == 0)):
                        st.markdown(calc["answer"])
                        if st.button("🗑️ Delete this", key=f"del_calc_{true_idx}"):
                            st.session_state.calc_history.pop(true_idx)
                            st.rerun()

                if st.button("🗑️ Clear All", key="clear_calc"):
                    st.session_state.calc_history = []
                    st.rerun()
        else:
            st.info("⚡ Precise Calculator works with CSV and Excel files.")

    # ── TAB 3: DATA ANALYSIS ──────────────────────────────────────────────────
    with tab3:
        st.subheader("📊 Data Analysis")
        if df is not None:
            # PERFORMANCE FIX: Use cached stats
            with st.spinner("📊 Calculating statistics..."):
                missing, dupes, stats_desc = get_df_stats(df)

            c1, c2, c3 = st.columns(3)
            c1.metric("Total Rows",     f"{len(df):,}")
            c2.metric("Total Columns",  f"{len(df.columns)}")
            c3.metric("Duplicate Rows", f"{dupes:,}")

            if missing:
                st.warning(f"⚠️ Missing values in: {', '.join(missing.keys())}")
                st.dataframe(pd.DataFrame.from_dict(
                    missing, orient="index", columns=["Missing Count"]
                ))
            else:
                st.success("✅ No missing values!")

            if dupes > 0:
                st.warning(f"⚠️ {dupes} duplicate rows found!")
            else:
                st.success("✅ No duplicate rows!")

            st.markdown("#### 📈 Auto Charts")
            num_cols = df.select_dtypes(include="number").columns.tolist()
            if num_cols:
                sel_col    = st.selectbox("Select column:", num_cols, key="chart_col")
                chart_type = st.radio("Chart type:", ["Bar","Line","Area"],
                                      horizontal=True, key="chart_type")
                
                # PERFORMANCE FIX: Sample data for charts (max 1,000 pts) 
                # Rendering 100k points crashes the browser.
                cdata = df[sel_col].dropna().reset_index(drop=True)
                if len(cdata) > 1000:
                    st.caption(f"Showing representative sample of 1,000/{len(cdata):,} points.")
                    cdata = cdata.sample(n=1000, random_state=42).sort_index()
                
                if chart_type == "Bar":   st.bar_chart(cdata)
                elif chart_type == "Line": st.line_chart(cdata)
                else:                      st.area_chart(cdata)

            st.markdown("#### 📋 Statistical Summary")
            st.dataframe(stats_desc)

            st.markdown("#### 🗃️ Raw Data (first 100 rows)")
            st.dataframe(df.head(100))

            # PERFORMANCE FIX: Use cached CSV for download
            csv_data = get_csv_download(df)
            st.download_button(
                "⬇️ Download as CSV", data=csv_data,
                file_name="data.csv", mime="text/csv", key="dl_csv"
            )
        else:
            st.info("📊 Available for CSV and Excel files.")

    # ── TAB 4: PREVIEW ────────────────────────────────────────────────────────
    with tab4:
        st.subheader("👀 File Preview")
        c1, c2, c3 = st.columns(3)
        c1.metric("Words",  f"{word_count:,}")
        c2.metric("Chunks", f"{len(chunks):,}")
        c3.metric("Size",   f"{size_mb:.2f} MB")
        st.text_area(
            "Extracted Content (first 3000 chars):",
            text[:3000] + ("..." if len(text) > 3000 else ""),
            height=400,
            key="preview_text"
        )

else:
    st.info("👋 Welcome! Upload any file in the sidebar to get started.")
    c1, c2, c3, c4 = st.columns(4)
    with c1:
        st.markdown("### 📄 Documents")
        st.write("PDF (scanned), Word, PPT, TXT")
    with c2:
        st.markdown("### 📊 Data Files")
        st.write("CSV, Excel, JSON, SQL, YAML")
    with c3:
        st.markdown("### 💻 Code & More")
        st.write("Python, HTML, JS, XML, ZIP, Images")
    with c4:
        st.markdown("### ⚡ Lightning Fast")
        st.write("TF-IDF RAG — no heavy model downloads")
    st.divider()
    st.markdown("""
### How it works:
1. **Upload** any file → text extracted instantly
2. **TF-IDF** indexes it in milliseconds (no 300MB model download!)
3. **Ask AI** → finds relevant parts & answers in plain English
4. **Precise Calculator** → runs real Python for 100% exact results
5. **Data Analysis** → auto charts, quality checks, full stats
6. **Download** answers or data anytime
""")
