import streamlit as st
import pandas as pd
import numpy as np
import io
import os
import re
import json
import time
import zipfile
from groq import Groq
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# ── CONFIG ────────────────────────────────────────────────────────────────────
MAX_FILE_MB    = 1024
MAX_RAG_ROWS   = 2000
MAX_WORDS      = 30000
TOP_K          = 3
CONTEXT_CAP    = 4000
MAX_OUT_TOKENS = 600

# ── API SETUP ─────────────────────────────────────────────────────────────────
def get_client():
    api_key = st.secrets.get("GROQ_API_KEY") or os.getenv("GROQ_API_KEY")
    if api_key:
        return Groq(api_key=api_key)
    return None

client = get_client()
has_api = client is not None

def check_api():
    if not has_api:
        st.error("🔑 GROQ_API_KEY missing. Please add it to .env or Streamlit Secrets.")
        return False
    return True

# ── PREMIUM UI ENGINE ──────────────────────────────────────────────────────────
def inject_lumina_css():
    """Injects the ultra-premium 'Chrome Abstract' design system from Stitch."""
    st.markdown("""
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Space+Grotesk:wght@300;400;500;600;700&family=Outfit:wght@300;400;600;800&family=Inter:wght@400;600;800&display=swap');
        
        :root {
            --primary: #bb00ff;
            --primary-glow: rgba(187, 0, 255, 0.4);
            --background: #0a0a0c;
            --glass: rgba(255, 255, 255, 0.02);
            --glass-border: rgba(187, 0, 255, 0.15);
            --chrome-gradient: linear-gradient(135deg, #FFFFFF 0%, #bb00ff 50%, #7a00ff 100%);
        }

        * { font-family: 'Outfit', 'Inter', sans-serif; }
        
        /* Global App Override */
        .stApp {
            background: radial-gradient(circle at top left, #1a0b2e 0%, #0f172a 45%, #0a0a0c 100%) !important;
            color: #E0E0E0 !important;
        }

        /* Hide Streamlit elements for lean look */
        header, footer {visibility: hidden;}
        #MainMenu {visibility: hidden;}

        /* Hero Text - Chrome Abstract Style */
        .hero-text {
            font-family: 'Space Grotesk', sans-serif;
            background: var(--chrome-gradient);
            background-size: 200% auto;
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
            font-weight: 800;
            font-size: 5rem !important;
            letter-spacing: -3px;
            margin-bottom: 0px;
            animation: chrome-shine 8s linear infinite;
            filter: drop-shadow(0 0 20px rgba(187, 0, 255, 0.2));
            line-height: 1.1;
            text-transform: uppercase;
        }

        @keyframes chrome-shine {
            to { background-position: 200% center; }
        }

        .lumina-subtext {
            font-family: 'Space Grotesk', sans-serif;
            color: #888;
            font-size: 1.2rem;
            font-weight: 300;
            letter-spacing: 5px;
            text-transform: uppercase;
            margin-top: -10px;
            opacity: 0.8;
        }

        /* Premium Glass Cards */
        .lumina-card, .glass-card {
            background: var(--glass);
            border: 1px solid var(--glass-border);
            border-radius: 24px;
            padding: 2.5rem;
            backdrop-filter: blur(15px);
            -webkit-backdrop-filter: blur(15px);
            transition: all 0.4s cubic-bezier(0.4, 0, 0.2, 1);
            margin-bottom: 25px;
            position: relative;
            overflow: hidden;
        }
        
        .lumina-card::before {
            content: '';
            position: absolute;
            top: 0; left: -100%;
            width: 100%; height: 100%;
            background: linear-gradient(90deg, transparent, rgba(187, 0, 255, 0.05), transparent);
            transition: 0.5s;
        }

        .lumina-card:hover::before {
            left: 100%;
        }

        .lumina-card:hover {
            border: 1px solid rgba(187, 0, 255, 0.4);
            background: rgba(187, 0, 255, 0.04);
            box-shadow: 0 0 30px rgba(187, 0, 255, 0.15);
            transform: translateY(-5px);
        }

        /* Neon Buttons */
        .stButton>button {
            background: linear-gradient(135deg, #bb00ff 0%, #7a00ff 100%) !important;
            color: white !important;
            border: none !important;
            border-radius: 14px !important;
            padding: 14px 28px !important;
            font-weight: 700 !important;
            font-family: 'Space Grotesk', sans-serif !important;
            text-transform: uppercase !important;
            letter-spacing: 1px !important;
            box-shadow: 0 8px 15px rgba(187, 0, 255, 0.2) !important;
            transition: all 0.3s ease !important;
        }
        .stButton>button:hover {
            transform: scale(1.02) !important;
            box-shadow: 0 0 25px rgba(187, 0, 255, 0.5) !important;
        }

        /* Status Badges */
        .badge-validated {
            padding: 6px 14px;
            border-radius: 999px;
            background: rgba(187, 0, 255, 0.1);
            color: #bb00ff;
            font-size: 11px;
            font-weight: 800;
            text-transform: uppercase;
            border: 1px solid rgba(187, 0, 255, 0.3);
            letter-spacing: 1px;
            box-shadow: 0 0 10px rgba(187, 0, 255, 0.2);
        }

        /* Chat UI Refinement */
        .ai-bubble {
            background: rgba(255, 255, 255, 0.02);
            backdrop-filter: blur(10px);
            border: 1px solid rgba(187, 0, 255, 0.1);
            border-radius: 0 24px 24px 24px;
            padding: 24px;
            margin-bottom: 25px;
            color: #ddd;
            font-size: 1rem;
            line-height: 1.6;
            box-shadow: 0 10px 30px rgba(0,0,0,0.3);
        }

        .user-bubble {
            background: linear-gradient(135deg, rgba(187, 0, 255, 0.1) 0%, rgba(122, 0, 255, 0.1) 100%);
            backdrop-filter: blur(10px);
            border: 1px solid rgba(187, 0, 255, 0.3);
            border-radius: 24px 24px 0 24px;
            padding: 24px;
            margin-bottom: 25px;
            margin-left: auto;
            max-width: 85%;
            color: #fff;
            box-shadow: 0 10px 30px rgba(187, 0, 255, 0.1);
        }

        /* Uploader Styles */
        [data-testid="stFileUploader"] {
            background: var(--glass);
            border: 1px dashed var(--glass-border);
            border-radius: 20px;
            padding: 20px;
        }
    </style>
    """, unsafe_allow_html=True)

def find_col(df, keywords):
    """Helper to find the best matching column for a list of keywords."""
    for col in df.columns:
        if any(k in col.lower() for k in keywords):
            return col
    return None

# ── UTILS ─────────────────────────────────────────────────────────────────────
@st.cache_data(show_spinner=False)
def get_df_stats(df):
    missing = df.isnull().sum()
    missing = missing[missing > 0].to_dict()
    dupes   = df.duplicated().sum()
    stats   = df.describe(include='all').transpose()
    return missing, dupes, stats

@st.cache_data(show_spinner=False)
def get_csv_download(df):
    return df.to_csv(index=False).encode('utf-8')

# ── EXTRACTORS ───────────────────────────────────────────────────────────────
@st.cache_data(show_spinner=False)
def extract_text_from_pdf(file_bytes):
    import pypdf
    reader = pypdf.PdfReader(io.BytesIO(file_bytes))
    text, ocr_pages = "", []
    for i, page in enumerate(reader.pages):
        pt = page.extract_text()
        if pt and len(pt.strip()) > 30:
            text += f"\n--- Page {i+1} ---\n{pt}\n"
        else:
            ocr_pages.append(i + 1)
    
    if ocr_pages:
        try:
            import pytesseract
            from pdf2image import convert_from_bytes
            imgs = convert_from_bytes(file_bytes, dpi=120)
            for pn in ocr_pages:
                if pn <= len(imgs):
                    ot = pytesseract.image_to_string(imgs[pn - 1])
                    if ot.strip():
                        text += f"\n--- Page {pn} (OCR) ---\n{ot}\n"
        except Exception:
            pass
    return text, None

@st.cache_data(show_spinner=False)
def extract_text_from_csv(file_bytes):
    try:
        df = pd.read_csv(io.BytesIO(file_bytes), low_memory=True)
        for col in df.select_dtypes(include=['float']).columns:
            df[col] = pd.to_numeric(df[col], downcast='float')
        for col in df.select_dtypes(include=['integer']).columns:
            df[col] = pd.to_numeric(df[col], downcast='integer')
        total = len(df)
        text  = f"CSV | Rows: {total:,} | Columns: {list(df.columns)}\n\n"
        if total > MAX_RAG_ROWS:
            sample = df.sample(n=MAX_RAG_ROWS, random_state=42)
            text += f"(Representative sample of {MAX_RAG_ROWS:,}/{total:,} rows for AI context.)\n\n"
            text += sample.to_string(index=False, max_rows=MAX_RAG_ROWS)
        else:
            text += df.to_string(index=False)
        return text, df
    except Exception as e:
        return f"Error parsing CSV: {e}", None

@st.cache_data(show_spinner=False)
def extract_text_from_excel(file_bytes):
    try:
        xl  = pd.ExcelFile(io.BytesIO(file_bytes), engine='openpyxl')
        txt, fdf = "", None
        for sheet in xl.sheet_names:
            df  = xl.parse(sheet)
            for col in df.select_dtypes(include=['float']).columns:
                df[col] = pd.to_numeric(df[col], downcast='float')
            if fdf is None: fdf = df
            txt += f"\n--- Sheet: {sheet} ---\nRows:{len(df)} Cols:{list(df.columns)}\n"
            txt += df.head(100).to_string(index=False) + "\n"
        return txt, fdf
    except Exception as e:
        return f"Excel Error: {e}", None

@st.cache_data(show_spinner=False)
def extract_text_from_docx(file_bytes):
    import docx
    doc  = docx.Document(io.BytesIO(file_bytes))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    for tbl in doc.tables:
        for row in tbl.rows:
            text += " | ".join(c.text.strip() for c in row.cells) + "\n"
    return text, None

@st.cache_data(show_spinner=False)
def extract_text_from_pptx(file_bytes):
    from pptx import Presentation
    prs  = Presentation(io.BytesIO(file_bytes))
    text = f"PowerPoint | Slides: {len(prs.slides)}\n"
    for i, slide in enumerate(prs.slides):
        text += f"\n--- Slide {i+1} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
    return text, None

@st.cache_data(show_spinner=False)
def extract_text_from_image(file_bytes):
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(io.BytesIO(file_bytes))
        txt = pytesseract.image_to_string(img)
        w, h = img.size
        return f"Image {w}x{h}px\n\n{txt if txt.strip() else 'No readable text.'}", None
    except Exception as e:
        return f"Image error: {e}", None

@st.cache_data(show_spinner=False)
def extract_text_from_zip(file_bytes):
    text = "ZIP Archive:\n"
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
        for name in z.namelist():
            text += f"  - {name}\n"
            if name.endswith((".txt", ".csv", ".json", ".md", ".py", ".sql")):
                try:
                    with z.open(name) as f:
                        raw = f.read().decode("utf-8", errors="ignore")
                        text += f"\n--- {name} ---\n"
                        text += (raw[:2000] + "...") if len(raw) > 2000 else raw
                        text += "\n"
                except Exception: pass
    return text, None

def extract_content(uploaded_file):
    name       = uploaded_file.name.lower()
    file_bytes = uploaded_file.getvalue()
    if name.endswith(".csv"): return extract_text_from_csv(file_bytes)
    if name.endswith((".xlsx", ".xls")): return extract_text_from_excel(file_bytes)
    if name.endswith(".pdf"): return extract_text_from_pdf(file_bytes)
    if name.endswith(".docx"): return extract_text_from_docx(file_bytes)
    if name.endswith(".pptx"): return extract_text_from_pptx(file_bytes)
    if name.endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff")): return extract_text_from_image(file_bytes)
    if name.endswith(".json"):
        try:
            data = json.loads(file_bytes.decode("utf-8", errors="ignore"))
            text = json.dumps(data, indent=2)[:8000]
            df = pd.DataFrame(data) if isinstance(data, list) and data and isinstance(data[0], dict) else None
            return text, df
        except Exception: return "Error reading JSON", None
    if name.endswith(".zip"): return extract_text_from_zip(file_bytes)
    return file_bytes.decode("utf-8", errors="ignore"), None

# ── AI LOGIC ──────────────────────────────────────────────────────────────────
@st.cache_data(show_spinner=False)
def generate_suggestions(text, file_name, is_tabular=False):
    if not check_api(): return ["Summary", "Trends", "Structure"]
    prompt = f"Based on '{file_name}', suggest 5 data-driven questions.\n\nContent:\n{text[:4000]}"
    try:
        r = client.chat.completions.create(model="llama-3.3-70b-versatile", messages=[{"role":"user","content":prompt}], max_tokens=300)
        return [l.strip().lstrip("0123456789.- ") for l in r.choices[0].message.content.split("\n") if l.strip()]
    except Exception: return ["Detailed summary", "Key trends", "Data structure"]

@st.cache_resource
def build_tfidf_index(text):
    if not text: return None, None, []
    words = text.split()[:MAX_WORDS]
    chunks, size, overlap = [], 400, 100
    for i in range(0, len(words), size - overlap):
        c = " ".join(words[i:i + size])
        if c.strip(): chunks.append(c)
    if not chunks: return None, None, []
    v = TfidfVectorizer(max_features=2000, stop_words="english")
    m = v.fit_transform(chunks)
    return v, m, chunks

def retrieve_chunks(question, vectorizer, tfidf_matrix, chunks):
    if not vectorizer or tfidf_matrix is None or not chunks:
        return chunks[:TOP_K] if chunks else []
    try:
        q_vec = vectorizer.transform([question])
        sims = cosine_similarity(q_vec, tfidf_matrix).flatten()
        top_indices = sims.argsort()[-TOP_K:][::-1]
        results = [chunks[i] for i in top_indices if sims[i] > 0]
        return results if results else chunks[:1]
    except Exception:
        return chunks[:1]

def ask_ai(question, chunks, history):
    if not check_api(): return "API Key Missing"
    context = "\n---\n".join(chunks)[:CONTEXT_CAP]
    messages = [{"role": "system", "content": "Assistant. Answer using ONLY provided text."}]
    if history:
        for c in history[-3:]:
            messages.append({"role": "user", "content": c["question"]})
            messages.append({"role": "assistant", "content": c.get("answer", "")})
    messages.append({"role": "user", "content": f"Context:\n{context}\n\nQ: {question}"})
    try:
        r = client.chat.completions.create(model="llama-3.3-70b-versatile", messages=messages, max_tokens=MAX_OUT_TOKENS)
        return r.choices[0].message.content.strip()
    except Exception as e: return f"Error: {e}"

def run_code(question, df, history, execute=True):
    if not check_api(): return "API Key Missing", "Error"
    cols = str(list(df.columns))[:1000]
    messages = [
        {"role": "system", "content": "You are a professional Python Data Analyst. Return ONLY the raw code required for the analysis. DO NOT explain. DO NOT use markdown code blocks. Use the dataframe 'df'. Important: Ensure all strings are properly closed and use single lines where possible. Output to Streamlit using 'st.metric' or 'st.write'."}
    ]
    if history:
        for c in history[-2:]:
            messages.append({"role": "user", "content": c["question"]})
            messages.append({"role": "assistant", "content": c.get("answer", "")})
    messages.append({"role": "user", "content": f"Data Columns: {cols}\nRequest: {question}"})
    
    raw = ""
    try:
        r = client.chat.completions.create(model="llama-3.3-70b-versatile", messages=messages, max_tokens=600)
        raw = r.choices[0].message.content.replace("```python", "").replace("```", "").strip()
        
        # Basic Security
        dangerous = ["os.", "sys.", "subprocess", "open(", "eval(", "exec(", "shutil"]
        if any(t in raw for t in dangerous): return raw, "Security block"
        
        if execute:
            exec(raw, {"df": df, "st": st, "pd": pd, "np": np})
        return raw, None
    except Exception as e:
        return raw, str(e)

def auto_generate_pulse(df, chunks):
    if not has_api: return "• Ready."
    prompt = f"Provide 3 bullet observations from this data summary:\n{df.head(2).to_string() if df is not None else chunks[0][:1000]}"
    try:
        r = client.chat.completions.create(model="llama-3.3-70b-versatile", messages=[{"role":"user","content":prompt}], max_tokens=150)
        return r.choices[0].message.content.strip()
    except Exception: return "• Analysis ready."
